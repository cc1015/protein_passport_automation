from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pathlib import Path
from models.protein_model.human_protein import HumanProtein
from models.protein_model.ortholog import Ortholog
from dataclasses import dataclass, field
from models.image import Img

@dataclass
class Entry:
    """
    Represents a protein passport entry (pptx).

    Attributes:
        template_path (str): Path of the template ppt.
        human (HumanProtein): HumanProtein of this Entry.
        orthologs (list): List of Orthologs of this Entry.
        user_name (str): User's name.
        powerpoint (Presentation): Presentation object of this Entry.
        slides (Slides): Slides of this Entry's Presentation.
        table_cells (list): List containing text to fill table cells of first slide in this Entry.
        output_path (Path): Output path.
    """
    template_path: str
    human: HumanProtein
    orthologs: list
    user_name: str
    powerpoint: Presentation = field(init=False)
    slides: list = field(init=False)
    table_cells: list = field(init=False)
    output_path: Path = field(init=False)

    def __post_init__(self):
        """
        Post init method for Entry. Sets powerpoint, slides, output_path, and table_cells fields.
        """
        self.powerpoint = Presentation(self.template_path)
        self.slides = self.powerpoint.slides
        self.output_path = Path(__file__).parent.parent.parent / "output" / f"{self.human.name}_protein_passport.pptx"
        self._build_table_cells()
        self._set_footer()

    def _build_table_cells(self):
        """
        Sets table_cells field.
        """
        self.table_cells = [
            [
                "Target Name: " + self.human.passport_table_data["rec_name"],
                "Aliases: " + ", ".join(self.human.passport_table_data["aliases"]),
                f"(Gene id: {self.human.name}, UniProtKB - {self.human.id})"
            ],
            [self.human.passport_table_data["target_type"]],
            ["interactions"],
            [
                f"{self.human.passport_table_data['length']} aa {self.human.passport_table_data['mass']} kDa",
                'self.human.nature_info'
            ],
            [f"{o.organism.value[0]}: {o.ecd_similarity}" for o in self.orthologs if o.ecd_similarity is not None],
            [
                f"Experimental PDBs: {', '.join(self.human.passport_table_data['exp_pdbs'])}",
                f"Predicted: {self.human.pred_pdb_id}"
            ],
            [f"{self.human.passport_table_data['exp_pattern']}."],
            [f"{self.human.passport_table_data['known_activity']}."]
        ]

    def _set_footer(self):
        """
        Writes user's name to footer.
        """
        for slide in self.slides:
            for shape in slide.shapes:
                if 'Footer' in shape.name:
                    shape.text=self.user_name
                    break

    def populate_info_table_slide(self, img: Img):
        """
        Populates the first slide of protein passport ppt template.

        Args:
            img (Img): Human protein 3d structure image.
        """
        slide = self.slides[0]

        for shape in slide.shapes:
            if 'Title' in shape.name:
                title = shape
            if 'Picture' in shape.name:
                picture = shape
            if 'Text' in shape.name:
                pbd_id_caption = shape
            if shape.has_table:
                table = shape.table
        
        if title:
            title.text = "Protein Passport - " + self.human.name

        if picture:
            img.vertical()
            picture.insert_picture(img.path)
        
        if pbd_id_caption:
            pbd_id_caption.text = img.caption
            pbd_id_caption.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        
        if table:
            for i, cell_text in enumerate(self.table_cells):
                cell = table.cell(i, 1)
                cell.text = "\n".join(cell_text)

                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        run.font.bold = False
                        run.font.color.rbg = RGBColor(0, 0, 0)

        self.powerpoint.save(self.output_path)
    
    def populate_str_align_slide(self, align_imgs: list, seq_img: Img=None):
        """
        Populates the third slide of protein passport ppt template.

        Args:
            align_imgs (list): Structure align images.
            seq_img (Img): Aligned sequence image.
        """
        slide = self.slides[2]
        img_1 = align_imgs[0]
        img_2 = align_imgs[1]

        #seq_img.horizontal()
        img_1.vertical()
        img_2.vertical()

        pictures = [img_1.path, img_2.path]
        placeholders = []
        captions = [img_1.caption, img_2.caption]
        textboxes = []

        for shape in slide.shapes:
            if 'Picture' in shape.name:
                placeholders.append(shape)
            if 'Table' in shape.name:
                table = shape.table
            if 'TextBox' in shape.name and len(textboxes) < 2:
                textboxes.append(shape)
        
        zipped = zip(placeholders[1:], pictures)
        for z in zipped:
            z[0].insert_picture(z[1])

        zipped = zip(textboxes, captions)
        for z in zipped:
            z[0].text = z[1]
            for paragraph in z[0].text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)

        if table:
            cell = table.cell(1, 1)
            cell.text = self.human.id
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

            for i, ortholog in enumerate(self.orthologs, start=1):
                species_cell = table.cell(i, 0)
                species_cell.text = ortholog.organism.name.capitalize()
                species_cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

                id_cell = table.cell(i, 1)
                id_cell.text = ortholog.id
                id_cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

                similarity_cell = table.cell(i, 2)
                similarity_cell.text = str(ortholog.ecd_similarity) + "%"
                similarity_cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

        self.powerpoint.save(self.output_path)
    
    def populate_string_db_slide(self, network_img: str, pred_partners_img=None):
        """
        Populates the fourth slide of protein passport ppt template.

        Args:
            network_img (str): path to STRING network image.
            pred_partners_img (str): path to STRING predicted partners.
        """
        slide = self.slides[3]
        placeholders = []

        for shape in slide.shapes:
            if 'Picture' in shape.name:
                placeholders.append(shape)
        
        placeholders[0].insert_picture(network_img)
        # placeholders[1].insert_picture(pred_partners_img)

        self.powerpoint.save(self.output_path)
    
    